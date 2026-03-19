from io import BytesIO

from fastapi import APIRouter, HTTPException
from fastapi.responses import Response
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches

router = APIRouter(prefix="/api", tags=["export"])


class ExportSlide(BaseModel):
    title: str
    subtitle: str | None = None
    bullets: list[str] = []
    layout: str = "bullet-list"


class ExportRequest(BaseModel):
    deckTitle: str = "Presentation"
    slides: list[ExportSlide]


def _add_slide(prs: Presentation, slide_data: ExportSlide) -> None:
    has_bullets = bool(slide_data.bullets)
    layout_idx = 1 if has_bullets else 0  # 0=title, 1=title+content
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes

    if shapes.title:
        shapes.title.text = slide_data.title

    if len(shapes.placeholders) < 2:
        return

    body = shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()

    lines: list[str] = []
    if slide_data.subtitle:
        lines.append(slide_data.subtitle)
    lines.extend(slide_data.bullets)

    if lines:
        tf.text = lines[0]
        for line in lines[1:]:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0


@router.post("/export/pptx")
async def export_pptx(payload: ExportRequest) -> Response:
    if not payload.slides:
        raise HTTPException(status_code=400, detail="No slides to export")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    for s in payload.slides:
        _add_slide(prs, s)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    filename = f"{payload.deckTitle or 'presentation'}.pptx".replace(" ", "_")
    return Response(
        content=buffer.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )
